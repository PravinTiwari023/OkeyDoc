import io
import os
import tempfile
from pdf2image import convert_from_bytes
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
import fitz

def ppt_to_pdf_service(input_bytes, quality='high'):
    """
    Convert PowerPoint document to PDF with specified quality
    
    Args:
        input_bytes: The input PowerPoint file as bytes
        quality: Quality level ('low', 'medium', 'high')
        
    Returns:
        PDF file as bytes
    """
    try:
        # Create temporary files for input and output
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_input:
            temp_input.write(input_bytes)
            temp_input_path = temp_input.name

        temp_output_path = os.path.splitext(temp_input_path)[0] + '.pdf'

        try:
            # Load presentation using python-pptx
            presentation = Presentation(temp_input_path)
            
            # Convert each slide to image
            images = []
            for slide in presentation.slides:
                # Save slide as temporary image
                img_path = f"{temp_input_path}_slide_{len(images)}.png"
                with open(img_path, 'wb') as img_file:
                    # Use pdf2image to convert slide to image
                    # This is a workaround since python-pptx doesn't support direct PDF export
                    slide_bytes = io.BytesIO()
                    presentation.save(slide_bytes)
                    slide_bytes.seek(0)
                    img = convert_from_bytes(slide_bytes.getvalue())[0]
                    img.save(img_path, 'PNG')
                images.append(img_path)

            # Create PDF from images
            images_pil = [Image.open(img_path) for img_path in images]
            if images_pil:
                images_pil[0].save(
                    temp_output_path,
                    'PDF',
                    save_all=True,
                    append_images=images_pil[1:]
                )

            # Read the output PDF
            with open(temp_output_path, 'rb') as pdf_file:
                pdf_bytes = pdf_file.read()

            return pdf_bytes

        finally:
            # Clean up temporary files
            try:
                os.remove(temp_input_path)
                os.remove(temp_output_path)
                for img_path in images:
                    if os.path.exists(img_path):
                        os.remove(img_path)
            except:
                pass

    except Exception as e:
        raise Exception(f'PowerPoint to PDF conversion failed: {str(e)}')

def pdf_to_ppt_service(input_bytes, quality='high'):
    """
    Convert PDF to PowerPoint with specified quality
    
    Args:
        input_bytes: The input PDF file as bytes
        quality: Quality level ('low', 'medium', 'high')
        
    Returns:
        PowerPoint file as bytes
    """
    try:
        # Quality settings
        dpi_settings = {
            'low': 150,
            'medium': 200,
            'high': 300
        }
        dpi = dpi_settings.get(quality, dpi_settings['high'])
        zoom = dpi / 72  # Convert DPI to zoom factor
        
        # Open PDF with PyMuPDF
        pdf_document = fitz.open(stream=input_bytes, filetype="pdf")
        prs = Presentation()
        
        # Process each page
        for page_num in range(len(pdf_document)):
            page = pdf_document[page_num]
            
            # Convert page to image
            pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            
            # Save image to bytes
            img_byte_arr = io.BytesIO()
            img.save(img_byte_arr, format='PNG', optimize=True)
            img_byte_arr.seek(0)
            
            # Add a slide with the image
            blank_slide_layout = prs.slide_layouts[6]  # blank layout
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Calculate image size to fit slide
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            img_width = Inches(pix.width / dpi * 72 / 25.4)
            img_height = Inches(pix.height / dpi * 72 / 25.4)
            
            # Scale image to fit slide while maintaining aspect ratio
            if img_width > slide_width or img_height > slide_height:
                scale = min(slide_width/img_width, slide_height/img_height)
                img_width *= scale
                img_height *= scale
            
            # Center image on slide
            left = (slide_width - img_width) / 2
            top = (slide_height - img_height) / 2
            
            # Add picture to slide
            slide.shapes.add_picture(img_byte_arr, left, top, width=img_width, height=img_height)
        
        # Save presentation to bytes
        pptx_bytes = io.BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)
        
        pdf_document.close()
        return pptx_bytes.getvalue()
        
    except Exception as e:
        raise Exception(f"PDF to PowerPoint conversion failed: {str(e)}")